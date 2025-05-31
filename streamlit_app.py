
# streamlit_app.py

prompt_library = {
    "custom": ""
}

system_message = """
You are a Tally Definition Language (TDL) expert and educator with deep experience in Tally Prime 6.0. The user will provide a task in plain English, and you will:

🔹 Refer to the sample TDL code snippets provided as your only source of truth.  
🔹 Generate a **complete, beginner-friendly instructional guide** that helps the user implement the requested TDL functionality step-by-step.  
🔹 Assume the user has **no prior knowledge of TDL or programming**.  
🔹 Your job is to **teach**.

---

** Note: Ensure you provide Snippets from Actual code samples in your guide. **

## 🧠 Reference Material:
Below you are given a set of real-world TDL code snippets. These come from live production code and represent best practices. You must learn from them to understand TDL syntax, naming conventions, object structure, and formatting.
    """

import os
import json
from typing import List, Dict, Tuple
import re
import requests # type: ignore
import tempfile
import pickle
import fitz  # PyMuPDF # type: ignore
import faiss # type: ignore
import streamlit as st # type: ignore
from sentence_transformers import SentenceTransformer # type: ignore
import numpy as np # type: ignore
from datetime import datetime, timedelta
import pytz # type: ignore
import time 
import boto3 # type: ignore
from tavily import TavilyClient # type: ignore
# import html
from streamlit_cookies_manager import EncryptedCookieManager # type: ignore
import streamlit_authenticator as stauth # type: ignore
import matplotlib.pyplot as plt # type: ignore
import plotly.express as px # type: ignore
import plotly.graph_objects as go # type: ignore
from docx import Document # type: ignore
from pptx import Presentation # type: ignore
import pandas as pd # type: ignore
# from azure.storage.blob import BlobServiceClient # type: ignore
import io
from io import BytesIO
from PIL import Image # type: ignore
import cv2 # type: ignore
import smtplib
# from email.mime.multipart import MIMEMultipart
# from email.mime.text import MIMEText
# from openai import AzureOpenAI, OpenAI # type: ignore
import copy
import uuid
import tiktoken # type: ignore
ENC = tiktoken.get_encoding("cl100k_base")  # ~15 % error on Claude 3
CTX_LIMIT = 200_000

def n_tokens(text: str) -> int:
    return len(ENC.encode(text))

def used_tokens(sys_msg, rag_context):
    total  = len(ENC.encode(sys_msg))
    total += len(ENC.encode(rag_context))
    return total

if "Authenticator" not in st.session_state:
    st.session_state["Authenticator"] = None
if "authenticated" not in st.session_state:
    st.session_state["authenticated"] = False
if "username" not in st.session_state:
    st.session_state["username"] = ""
if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False
if "logged_out" not in st.session_state:
    st.session_state["logged_out"] = False


# Instantiate the Cookie Manager at the very top.
cookies = EncryptedCookieManager(
    prefix="tdl_copilot/",  # Use a unique prefix for your app.
    password="nexusdms"
)

if not cookies.ready():
    st.spinner("Loading cookies...")
    st.stop()

def remove_last_line(text):
    lines = text.splitlines()
    filtered_lines = [
        line for line in lines 
        if 'python' not in line.lower() and 'plotly' not in line.lower() and "code" not in line.lower()
    ]
    return "\n".join(filtered_lines)

# Define valid users in a dictionary
def load_dict_from_json(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        data = json.load(file)
    return data

# Load the MPNet model
mpnet_model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")

# Replace these with your actual AWS credentials
REGION = "ap-south-1"
TAVILY_API = "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"

# Paths for saving index and metadata
FAISS_INDEX_PATH = "faiss_index.bin"
METADATA_STORE_PATH = "metadata_store.pkl"

# AWS S3 setup
s3_bucket_name = "tdl-copilot-v1"

# Define a helper function to display your company logo
def display_logo():
    # Make sure 'logo.png' is in your working directory
    # st.image("../logo.png", width=150)
    st.write("<Your Logo Here>")

# Create a Textract Runtime client for document analysis
textract_client = boto3.client('textract', region_name=REGION)

# Create an S3 client for storage
s3_client = boto3.client('s3', region_name=REGION)

def get_presigned_url(file_key, expiration=3600):
    """
    Generate a pre-signed URL for the S3 object.
    """
    return s3_client.generate_presigned_url(
        'get_object',
        Params={'Bucket': s3_bucket_name, 'Key': file_key},
        ExpiresIn=expiration
    )

def save_chat_history(chat_history, blob_name="chat_history.json"):
    try:
        local_file_path = "chat_history.json"
        with open(local_file_path, "w", encoding="utf-8") as f:
            json.dump(chat_history, f, ensure_ascii=False, indent=2)
        s3_client.upload_file(local_file_path, s3_bucket_name, blob_name)
        # os.remove(local_file_path)
    except Exception as e:
        st.error(f"Failed to save chat history: {str(e)}")

def load_chat_history(blob_name="chat_history.json"):
    try:
        response = s3_client.list_objects_v2(Bucket=s3_bucket_name, Prefix=blob_name)
        if 'Contents' in response:
            local_file_path = "chat_history.json"
            s3_client.download_file(s3_bucket_name, blob_name, local_file_path)
            chat_history = json.load(open(local_file_path, encoding="utf-8"))
            # Make sure the file is a dict
            if not isinstance(chat_history, dict):
                chat_history = {}
            # Optional: fix any inner structure if needed
            return chat_history
        return {}
    except Exception as e:
        st.error(f"Failed to load chat history: {str(e)}")
        return {}


def file_exists_in_blob(file_name):
    """Check if a file with the same name exists in S3."""
    try:
        s3_client.head_object(Bucket=s3_bucket_name, Key=file_name)
        return True
    except Exception as e:
        if e.response['Error']['Code'] == '404':
            return False
        else:
            raise e  # Re-raise other exceptions

# Function to upload file to Azure Blob Storage
def upload_to_blob_storage(local_file_path, bucket_name, s3_key):
    try:
        with open(local_file_path, "rb") as data:
            s3_client.upload_fileobj(data, bucket_name, s3_key)

        # st.success(f"File '{s3_key}' successfully uploaded to S3.")
    except Exception as e:
        st.error(f"Failed to upload file to S3: {str(e)}")

# Function to download file from Azure Blob Storage
def download_from_blob_storage(s3_bucket_name, s3_key, local_file_path):
    """Download a file from S3, or return False if not found."""
    try:
        with open(local_file_path, "wb") as file:
            s3_client.download_fileobj(s3_bucket_name, s3_key, file)
        return True
    except Exception as e:
        if e.response['Error']['Code'] == '404':
            print(f"File not found in S3: {s3_key}")
            return False
        else:
            print(f"Failed to download {s3_key}: {str(e)}")
            return False

def create_word_doc(text):
    doc = Document()
    doc.add_heading("Chat Answer", level=1)
    doc.add_paragraph(text)
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Function to Generate titan embeddings
def generate_titan_embeddings(text):
    try:
        # Generate embeddings using MPNet
        embedding = mpnet_model.encode(text, normalize_embeddings=True)
        return np.array(embedding)
    except Exception as e:
        print(f"Error generating embeddings: {e}")
        return None  # Return None to handle errors gracefully

def call_llm_api(system_message, user_query):
    """
    Invoke the Anthropic Claude 3 Sonnet model on Amazon Bedrock with a text prompt.

    Args:
        prompt (str): The user's text prompt to send to the model.
        model_id (str): The Bedrock model ID for Claude 3 Sonnet.
        max_tokens (int): Maximum number of tokens to generate.
        anthropic_version (str): The version tag for the Anthropic API.

    Returns:
        str: The generated text response from the model.

    Raises:
        ClientError: If the Bedrock invoke_model call fails.
    """
    prompt = system_message + user_query
    # Create Bedrock Runtime client
    client = boto3.client("bedrock-runtime", region_name=REGION)  # :contentReference[oaicite:0]{index=0}

    # Assemble the native Bedrock payload
    native_request = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": prompt
                    }
                ]
            }
        ]
    }

    try:
        # Invoke the model
        response = client.invoke_model(
            modelId="arn:aws:bedrock:ap-south-1:273354629305:inference-profile/apac.anthropic.claude-3-5-sonnet-20241022-v2:0",
            contentType="application/json",
            accept="application/json",
            body=json.dumps(native_request)
        )
        # The body is a JSON-encoded string
        body_str = response["body"].read().decode("utf-8")
        body = json.loads(body_str)
        return body['content'][0]['text']
    
    except Exception as e:
        return f"An error occurred: {str(e)}"

# Faiss index initialization
dimension = 768  # Embedding dimension for text embeddings v3
faiss_index = faiss.IndexFlatL2(dimension)
metadata_store = []

# Updated `save_index_and_metadata` to upload files to Azure Blob Storage
def save_index_and_metadata():
    # Save files locally
    faiss.write_index(faiss_index, FAISS_INDEX_PATH)
    with open(METADATA_STORE_PATH, "wb") as f:
        pickle.dump(metadata_store, f)

    # Upload files to Blob Storage
    try:
        upload_to_blob_storage(FAISS_INDEX_PATH, s3_bucket_name, os.path.basename(FAISS_INDEX_PATH))
        upload_to_blob_storage(METADATA_STORE_PATH, s3_bucket_name, os.path.basename(METADATA_STORE_PATH))
    except Exception as e:
        print(f"Error uploading index or metadata to Blob Storage: {str(e)}")

# Function to load index and metadata
# Load index and metadata from Azure Blob Storage or initialize new
def load_index_and_metadata():
    global faiss_index, metadata_store

    index_blob_name = os.path.basename(FAISS_INDEX_PATH)
    metadata_blob_name = os.path.basename(METADATA_STORE_PATH)

    # Download files from Blob Storage if available
    index_downloaded = download_from_blob_storage(s3_bucket_name, index_blob_name, FAISS_INDEX_PATH)
    metadata_downloaded = download_from_blob_storage(s3_bucket_name, metadata_blob_name, METADATA_STORE_PATH)

    if index_downloaded and metadata_downloaded:
        # Load FAISS index and metadata store
        try:
            faiss_index = faiss.read_index(FAISS_INDEX_PATH)
            with open(METADATA_STORE_PATH, "rb") as f:
                metadata_store = pickle.load(f)
            print("Index and metadata loaded from Storage.")
        except Exception as e:
            print(f"Failed to load index or metadata: {str(e)}")
            # Initialize empty index and metadata if loading fails
            faiss_index = faiss.IndexFlatL2(dimension)
            metadata_store = []
    else:
        print("Index or metadata not found in Blob Storage. Initializing new.")
        # Initialize empty index and metadata
        faiss_index = faiss.IndexFlatL2(dimension)
        metadata_store = []

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    pages_text = []
    processing_message_placeholder = st.empty()
    progress_bar = st.progress(0)

    for page_num in range(len(doc)):
        processing_message_placeholder.write(f"Processing page {page_num + 1}/{len(doc)}...")
        temp_image_path = None

        try:
            # Render page as an image
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            temp_image_path = os.path.join(tempfile.gettempdir(), f"page_{page_num}.png")
            pix.save(temp_image_path)

            with open(temp_image_path, "rb") as image_file:
                image_bytes = image_file.read()

            # Use basic OCR (DetectDocumentText) instead of full analysis
            response = textract_client.detect_document_text(
                Document={'Bytes': image_bytes}
            )

            page_content = {"page_num": page_num + 1, "text": ""}
            
            # Extract text directly from OCR results
            text_lines = []
            for block in response.get('Blocks', []):
                if block['BlockType'] == 'LINE' and 'Text' in block:
                    text_lines.append(block['Text'])
            
            page_content["text"] = "\n".join(text_lines)
            pages_text.append((page_num + 1, page_content['text'], page_content['text']))

        except Exception as e:
            st.error(f"Error processing page {page_num + 1}: {str(e)}")
            continue

        finally:
            if temp_image_path and os.path.exists(temp_image_path):
                os.remove(temp_image_path)

        progress_bar.progress((page_num + 1) / len(doc))

    return pages_text

def chunk_text(pages_text, chunk_size=1):
    """Create non-overlapping chunks of text from pages."""
    chunks = []
    total_pages = len(pages_text)

    # Loop to create non-overlapping chunks
    for i in range(0, total_pages, chunk_size):
        chunk_parts = []
        for j in range(chunk_size):
            if i + j < total_pages:  # Ensure we do not exceed total pages
                page_num, text = pages_text[i + j]
                chunk_parts.append(f"The following text is from Page {page_num}:\n``````\n\n")
        
        if chunk_parts:  # Only append non-empty chunks
            chunks.append(''.join(chunk_parts))

    return chunks

def delete_file(file_name):
    try:
        # Delete from Amazon S3:
        s3_client.delete_object(Bucket=s3_bucket_name, Key=file_name)
        st.success(f"Deleted file '{file_name}' from S3 bucket.")
    except Exception as e:
        st.error(f"Error deleting file: {str(e)}")
    
    # Remove the file from metadata_store (filter out all records with that filename)
    global metadata_store
    metadata_store = [md for md in metadata_store if md["filename"] != file_name]
    
    # Optionally, update your index and re-save metadata
    save_index_and_metadata()

def user_has_file_access(username, file_name):
    """
    Checks whether the given username already has access to a file with file_name.
    This means the file is present in metadata_store with the same filename 
    AND the user is either the owner or is in that file's shared_with list.
    """
    for record in metadata_store:
        if record.get("filename") == file_name:
            if record.get("owner") == username or username in record.get("shared_with", []):
                return True
    return False

def add_pdf_to_index(pdf_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        temp_pdf.write(pdf_file.read())
        temp_pdf_path = temp_pdf.name
        
        pages_text = extract_text_from_pdf(temp_pdf_path)
        
        total_pages = len(pages_text)
        progress_bar = st.progress(0)  # Start at 0%
        
        # Create a placeholder for processing messages
        processing_message_placeholder = st.empty()

        # page_clf_text = ""
        
        for page_num, text, page_content in pages_text:
            processing_message_placeholder.write(f"Processing page {page_num}/{total_pages}...")

            embedding = generate_titan_embeddings(text)
            faiss_index.add(embedding.reshape(1, -1))
            metadata_store.append({
                "filename": os.path.basename(pdf_file.name),
                "page": page_num,
                "text": page_content,
                "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
                "shared_with": [] 
            })            
            
            # Update progress bar and clear previous message after updating
            progress_bar.progress(page_num / total_pages)  # Update to reflect current page
        
        save_index_and_metadata()
        
        try:
            # Upload the processed file to Azure Blob Storage
            blob_name = os.path.basename(pdf_file.name)
            upload_to_blob_storage(temp_pdf_path, s3_bucket_name, blob_name)
            
            os.remove(temp_pdf_path)  # Ensure file is removed after processing
        except PermissionError:
            st.warning("Could not delete the temporary PDF file. It might still be in use.")

def get_page_range_for_files(selected_files):
    """
    Finds the min and max page numbers for the selected files from metadata.
    """
    page_ranges = {}
    
    for file in selected_files:
        # Filter metadata for the selected file
        file_metadata = [metadata for metadata in metadata_store if metadata['filename'] == file]
        
        # Get the minimum and maximum page numbers for the file
        if file_metadata:
            min_page = min(metadata['page'] for metadata in file_metadata)
            max_page = max(metadata['page'] for metadata in file_metadata)
            page_ranges[file] = (min_page, max_page)
    
    return page_ranges

###

def query_documents_with_page_range(
    selected_files,
    selected_page_range,
    user_message,
    top_k,
    last_messages,
    web_search
    ):
    """
    Refactored function to:
    1. Understand the user query.
    2. Fetch document context from the selected files.
    3. Perform web search using Tavily API for additional references.
    4. Combine document and web search contexts to frame a response.
    """

    # Step 1: Embed the user query and fetch top-K document contexts
    query_embedding = generate_titan_embeddings(user_message).reshape(1, -1)
    if faiss_index.ntotal == 0:
        st.error("The FAISS index is empty. Please upload a PDF to populate the index.")
        return "No data available to query."
    
    k = faiss_index.ntotal
    D, I = faiss_index.search(query_embedding, k)

    include_all = not selected_files    
    # 2) filter
    relevant = []
    for dist, idx in zip(D[0], I[0]):
        if idx < len(metadata_store):
            record = metadata_store[idx]

            # ---------- NEW LINE --------------------------------------------
            if include_all or record["filename"] in selected_files:
            # ----------------------------------------------------------------
                (start_pg, end_pg) = selected_page_range.get(
                    record["filename"], (1, 999_999)
                )
                if start_pg <= record["page"] <= end_pg:
                    relevant.append((dist, record))

    # sort ascending by distance, take top_k
    relevant_sorted = sorted(relevant, key=lambda x: x[0])[: top_k]
    top_k_metadata = [r[1] for r in relevant_sorted]

    # Step 3: Perform web search using Tavily API
    if web_search:
        try:
            client = TavilyClient(api_key=TAVILY_API)
            web_search_results = client.search(
                query=user_message,
                search_depth="advanced",
                include_raw_content=True
            )
        except Exception as e:
            st.error(f"Web search failed: {str(e)}")
            web_search_results = {}

    # Step 4: Combine contexts and query the LLM
    combined_context = f"""
    # User Query: {user_message}

    # Last Messages: {json.dumps(last_messages)}
    
    # Document Context:
    {json.dumps(top_k_metadata, indent=2)}
    """

    if web_search:
        combined_context += f"""
        # Web Search Results:
        {json.dumps(web_search_results, indent=2)}
        """

    # ── count before we call the LLM ──────────────────────────────
    tokens_used = used_tokens(system_message, combined_context)
    pct = min(1.0, tokens_used/CTX_LIMIT)
    st.sidebar.markdown(f"**Token use:** {tokens_used:,} / {CTX_LIMIT:,}")
    st.sidebar.progress(pct)
    if tokens_used > 0.95*CTX_LIMIT:
        st.warning("⚠️ Within 5 % of Skywork’s context window; answers may be clipped.")

    try:
        response = call_llm_api(
            system_message,
            combined_context
        )
        return response
    except Exception as e:
        st.error(f"LLM query failed: {str(e)}")
        return "An error occurred while querying the LLM."

###

# Function to handle user input with text_area and predefined prompts
def user_input():

    prompt_options = list(prompt_library.keys())
    selected_prompt = st.selectbox("Select a predefined prompt:", prompt_options, index=0)
    
    # Auto-fill text area if a prompt is selected
    default_text = prompt_library[selected_prompt]
    user_message = st.text_area("Enter your message:", value=default_text, height=150)
    ol1, ol2, ol3, ol4, ol5, ol6, ol7, ol8 = st.columns(8)
    with ol8:
    # Submit button
       submit = st.button("Send")
    if submit and user_message.strip():
        default_text = "custom"
        st.session_state.user_message = ""  # Clear the text input
        return user_message.strip()
    return None

USERS = {"user": "password"}
credentials = {"usernames": {}}
for user in USERS:
    # Assume each user has keys: 'username', 'name', and 'password'
    credentials["usernames"][user] = {
        "name": user,
        "password": USERS[user]
    }

def login():
    display_logo()
    st.title("Ready to Dive In? Sign In!")
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    
    if st.button("Login"):
        if username in USERS and password == USERS[username]:
            # Mark the user as authenticated
            st.session_state["authenticated"] = True
            st.session_state["username"] = username
            st.session_state["logged_in"] = True

            # Set a persistent cookie
            cookies["username"] = username
            cookies.save()
            
            # Create the authenticator instance and store it in session state.
            # (Ensure you have a 'credentials' variable ready.)
            Authenticator = stauth.Authenticate(credentials, cookie_name='nexusdms/', key='abcdefgh', cookie_expiry_days=0)
            st.session_state["Authenticator"] = Authenticator
            
            st.success("Login successful!")
            # Optionally, you can call st.rerun() here if needed.
        else:
            st.error("Invalid username or password.")

def logout():
    # Retrieve the stored authenticator instance
    authenticator = st.session_state.get("Authenticator")
    if authenticator is not None:
        try:
            # Attempt logout (this might raise a KeyError if cookie already removed)
            logout_button = authenticator.logout('Log Out', 'sidebar')
        except KeyError:
            logout_button = True  # If cookie already removed, treat as successful logout.
        except Exception as err:
            st.error(f'Unexpected exception during logout: {err}')
            return
    else:
        logout_button = True

    if logout_button:
        # Update session state to reflect logout
        st.session_state["logged_out"] = True
        st.session_state["logged_in"] = False
        st.session_state["authenticated"] = False
        st.session_state["username"] = ""
        st.session_state["Authenticator"] = None
        
        # Clear the cookie as well.
        if "username" in cookies:
            # del cookies["username"]
            cookies["username"] = ""
            cookies.save()
        st.rerun()

def process_pdf(pdf_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        temp_pdf.write(pdf_file.read())
        temp_pdf_path = temp_pdf.name
        
    pages_text = extract_text_from_pdf(temp_pdf_path)
    total_pages = len(pages_text)
    progress_bar = st.progress(0)
    processing_message_placeholder = st.empty()
    
    for page_num, text, page_content in pages_text:
        processing_message_placeholder.write(f"Processing page {page_num}/{total_pages}...")
        embedding = generate_titan_embeddings(text)
        faiss_index.add(embedding.reshape(1, -1))
        metadata_store.append({
            "filename": os.path.basename(pdf_file.name),
            "page": page_num,
            "text": page_content,
            "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
            "shared_with": [] 
        })
        progress_bar.progress(page_num / total_pages)
    
    save_index_and_metadata()
    try:
        blob_name = os.path.basename(pdf_file.name)
        upload_to_blob_storage(temp_pdf_path, s3_bucket_name, blob_name)
        os.remove(temp_pdf_path)
    except PermissionError:
        st.warning("Could not delete the temporary PDF file. It might still be in use.")

# ------------------------------------------------------------------------
# Handle .txt files  ➜  read ➜  embed ➜  index ➜  upload ➜  delete
# ------------------------------------------------------------------------
def process_text(text_file):
    """
    Reads a plain-text file, chunks it, embeds each chunk, stores metadata,
    then uploads the original file to S3 and removes the local temp copy.
    """
    import tempfile, os

    # 1️⃣  Read the UploadFile into memory
    content = text_file.read().decode("utf-8", errors="ignore")

    # 2️⃣  Chunk at ~1 000-char intervals so embeddings stay meaningful
    chunks = [content[i : i + 1_000] for i in range(0, len(content), 1_000)]

    for idx, chunk in enumerate(chunks, 1):
        embedding = generate_titan_embeddings(chunk)
        faiss_index.add(embedding.reshape(1, -1))
        metadata_store.append({
            "filename": text_file.name,
            "page": idx,                     # re-use “page” as a simple counter
            "text": chunk,
            "owner": st.session_state.get("username", "unknown"),
            "shared_with": []
        })

    save_index_and_metadata()

    # 3️⃣  Mirror the ORIGINAL file to S3, then delete the temp copy
    try:
        # Write the in-memory string to a temp file so our helper can read it
        with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as tmp:
            tmp.write(content.encode("utf-8"))
            temp_txt_path = tmp.name

        blob_name = os.path.basename(text_file.name)
        upload_to_blob_storage(temp_txt_path, s3_bucket_name, blob_name)
        os.remove(temp_txt_path)

    except PermissionError:
        st.warning("Could not delete the temporary text file—it might still be open.")


def extract_text_from_image(file_path):
    # Open image with fitz (PyMuPDF supports image reading too)
    img_doc = fitz.open(file_path)
    text_lines = []
    try:
        page = img_doc[0]  # Assume the entire image is one "page"
        pix = page.get_pixmap()
        temp_image_path = os.path.join(tempfile.gettempdir(), "temp_image.png")
        pix.save(temp_image_path)
        with open(temp_image_path, "rb") as image_file:
            image_bytes = image_file.read()
        response = textract_client.detect_document_text(Document={'Bytes': image_bytes})
        for block in response.get('Blocks', []):
            if block['BlockType'] == 'LINE' and 'Text' in block:
                text_lines.append(block['Text'])
        os.remove(temp_image_path)
    except Exception as e:
        st.error(f"Error processing image: {str(e)}")
    return "\n".join(text_lines)

def process_image(image_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_img:
        temp_img.write(image_file.read())
        temp_img_path = temp_img.name
    text = extract_text_from_image(temp_img_path)
    embedding = generate_titan_embeddings(text)
    faiss_index.add(embedding.reshape(1, -1))
    metadata_store.append({
        "filename": os.path.basename(image_file.name),
        "page": 1,  # Only one "page" for an image
        "text": text,
        "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
        "shared_with": [] 
    })
    save_index_and_metadata()
    try:
        blob_name = os.path.basename(image_file.name)
        upload_to_blob_storage(temp_img_path, s3_bucket_name, blob_name)
        os.remove(temp_img_path)
    except PermissionError:
        st.warning("Could not delete the temporary image file. It might still be in use.")

def process_docx(docx_file):
    # Save file locally
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_docx:
        temp_docx.write(docx_file.read())
        temp_docx_path = temp_docx.name

    try:
        doc = Document(temp_docx_path)
        # If page information is not available, you can split the document into chunks by paragraphs.
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip() != ""])
        # Optional: use a heuristic to approximate pages (e.g., split full_text into equal parts)
        chunks = [full_text[i:i+1000] for i in range(0, len(full_text), 1000)]
        for idx, chunk in enumerate(chunks, 1):
            embedding = generate_titan_embeddings(chunk)
            faiss_index.add(embedding.reshape(1, -1))
            metadata_store.append({
                "filename": os.path.basename(docx_file.name),
                "page": idx,
                "text": chunk,
                "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
                "shared_with": [] 
            })
    except Exception as e:
        st.error(f"Error processing DOCX: {str(e)}")
    finally:
        os.remove(temp_docx_path)
    save_index_and_metadata()

def process_pptx(pptx_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
        temp_pptx.write(pptx_file.read())
        temp_pptx_path = temp_pptx.name
    try:
        prs = Presentation(temp_pptx_path)
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            embedding = generate_titan_embeddings(slide_text)
            faiss_index.add(embedding.reshape(1, -1))
            metadata_store.append({
                "filename": os.path.basename(pptx_file.name),
                "page": idx,  # Slide number
                "text": slide_text,
                "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
                "shared_with": [] 
            })
    except Exception as e:
        st.error(f"Error processing PPTX: {str(e)}")
    finally:
        os.remove(temp_pptx_path)
    save_index_and_metadata()

def process_spreadsheet(file_obj):
    ext = os.path.splitext(file_obj.name)[1].lower()
    try:
        if ext == ".csv":
            df = pd.read_csv(file_obj)
            sheet_name = "csv"
        else:
            xls = pd.ExcelFile(file_obj)
            # Process all sheets or just the first one; here we use the first sheet
            sheet_name = xls.sheet_names[0]
            df = pd.read_excel(xls, sheet_name=sheet_name)
    except Exception as e:
        st.error(f"Error reading spreadsheet: {str(e)}")
        return

    chunk_size = 50
    num_chunks = (len(df) // chunk_size) + int(len(df) % chunk_size != 0)
    for i in range(num_chunks):
        chunk_df = df.iloc[i*chunk_size : (i+1)*chunk_size]
        chunk_text = chunk_df.to_string(index=False)
        embedding = generate_titan_embeddings(chunk_text)
        faiss_index.add(embedding.reshape(1, -1))
        metadata_store.append({
            "filename": f"{os.path.basename(file_obj.name)} ({sheet_name})",
            "page": i + 1,  # Store as an integer instead of "chunk_{i+1}"
            "text": chunk_text,
            "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
            "shared_with": [] 
        })
    save_index_and_metadata()

def add_file_to_index(uploaded_file):
    ext = os.path.splitext(uploaded_file.name)[1].lower()
    
    if ext == ".pdf":
        process_pdf(uploaded_file)
    elif ext in [".jpg", ".jpeg", ".png"]:
        process_image(uploaded_file)
    elif ext in [".doc", ".docx"]:
        process_docx(uploaded_file)
    elif ext == ".pptx":
        process_pptx(uploaded_file)
    elif ext in [".xlsx", ".csv"]:
        process_spreadsheet(uploaded_file)
    else:
        process_text(uploaded_file)

# --- Helper Function: Overlay Logo and Dynamic Text ---
def load_image_from_input(image_input):
    """
    Loads an image from a local file path or a URL.
    """
    if isinstance(image_input, str) and image_input.startswith("http"):
        try:
            response = requests.get(image_input)
            response.raise_for_status()  # Ensure we got a valid response
            return Image.open(BytesIO(response.content))
        except requests.RequestException as e:
            st.error(f"Error loading image from URL: {e}")
            return None
    else:
        return Image.open(image_input)

def overlay_logo_and_text(image_input, logo_path, bottom_text, overlay_text, overlay_x_pct, overlay_y_pct, text_mode, font_adjuster):
    # Load main image from URL or local file
    main_image = load_image_from_input(image_input)
    if main_image is None:
        return None  # Stop processing if image load failed

    # Load logo (assuming logo is a local file)
    logo = Image.open(logo_path)

    # Get image dimensions
    img_width, img_height = main_image.size

    # Resize logo relative to image size (10% of image width)
    logo_width = int(img_width * 0.1)
    logo_height = int(logo.size[1] * (logo_width / logo.size[0]))  # Maintain aspect ratio
    logo = logo.resize((logo_width, logo_height), Image.Resampling.LANCZOS)

    # Paste the logo at top-left with 5% padding
    logo_x = int(img_width * 0.05)
    logo_y = int(img_height * 0.05)
    main_image.paste(logo, (logo_x, logo_y), logo)

    # Convert image to OpenCV format for text overlay
    image_cv = np.array(main_image)
    image_cv = cv2.cvtColor(image_cv, cv2.COLOR_RGB2BGR)

    # Define font settings (scaling font and thickness with image width)
    font = cv2.FONT_HERSHEY_SIMPLEX
    font_scale = max(1, img_width // 500) * font_adjuster
    thickness = max(2, img_width // 400)

    # Set text color based on user selection
    if text_mode.lower() == "light":
        font_color = (255, 255, 255)  # White text for light backgrounds
    else:
        font_color = (0, 0, 0)        # Black text for dark backgrounds

    # Add bottom text at 5% padding from left and 5% above the bottom edge
    bottom_x = int(img_width * 0.05)
    bottom_y = img_height - int(img_height * 0.05)
    cv2.putText(image_cv, bottom_text, (bottom_x, bottom_y), font, font_scale, font_color, thickness, cv2.LINE_AA)

    # Calculate overlay text position using percentage values (0 to 100)
    overlay_x = int(img_width * overlay_x_pct / 100)
    overlay_y = int(img_height * overlay_y_pct / 100)
    cv2.putText(image_cv, overlay_text, (overlay_x, overlay_y), font, font_scale, font_color, thickness, cv2.LINE_AA)

    # Convert back to PIL format
    final_image = Image.fromarray(cv2.cvtColor(image_cv, cv2.COLOR_BGR2RGB))
    return final_image

def main():
    # Try to restore authentication from session state or cookie.
    if not st.session_state["authenticated"]:
        cookie_username = cookies.get("username")
        if cookie_username != "" and cookie_username is not None:
            st.session_state["authenticated"] = True
            st.session_state["username"] = cookie_username
            st.session_state["logged_in"] = True

    if not st.session_state["authenticated"]:
        login()
        return

    display_logo()
    st.title("Document Query Assistant")

    # Initialize session state variables if not present.
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'current_conversation' not in st.session_state:
        st.session_state.current_conversation = None
    if 'sources' not in st.session_state:
        st.session_state.sources = []
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = load_chat_history()
    # Initialize file selection state.
    if 'selected_files' not in st.session_state:
        st.session_state.selected_files = []
    if 'selected_page_ranges' not in st.session_state:
        st.session_state.selected_page_ranges = {}
    if 'file_summaries' not in st.session_state:
        st.session_state.file_summaries = {}
    if "rename_mode" not in st.session_state:
        st.session_state["rename_mode"] = None
    if "share_chat_mode" not in st.session_state:
        st.session_state["share_chat_mode"] = False
    if "share_chat_conv" not in st.session_state:
        st.session_state["share_chat_conv"] = None
    if "share_chat_conv_id" not in st.session_state:
        st.session_state["share_chat_conv_id"] = None

    available_usernames = list(USERS.keys())
    
    current_user = st.session_state["username"]
    load_index_and_metadata()
    st.sidebar.header(f"Hello `{current_user}`")
    if st.sidebar.button("Log Out"):
        logout()  # Display the logout button in the sidebar

    st.sidebar.header("Options")
    option = st.sidebar.selectbox("Choose an option", ["Query Documents", "Upload Documents", "File Manager", "Usage Monitoring"])

    if option == "Upload Documents":
        st.header("Upload Documents")
        uploaded_files = st.file_uploader(
            "Upload one or more documents",
            # type=["pdf", "jpg", "jpeg", "png", "doc", "docx", "pptx", "xlsx", "csv"],
            accept_multiple_files=True
        )

        if uploaded_files:
            for uploaded_file in uploaded_files:
                st.write(f"Processing {uploaded_file.name}...")
                current_user = st.session_state.get("username", "unknown")
                # if file_exists_in_blob(uploaded_file.name):
                #     st.warning(f"File '{uploaded_file.name}' already exists in Storage. Skipping upload.")
                # else:
                #     add_file_to_index(uploaded_file)
                #     st.success(f"File '{uploaded_file.name}' has been successfully uploaded and added to the index.")
                # Check if current user already has access
                if user_has_file_access(current_user, uploaded_file.name):
                    st.warning(f"File '{uploaded_file.name}' is already in your library (owned or shared). Skipping upload.")
                else:
                    # Proceed with normal upload & indexing
                    add_file_to_index(uploaded_file)
                    st.success(f"File '{uploaded_file.name}' has been successfully uploaded and added to your library.")

    elif option == "File Manager":
        # available_usernames = list(USERS.keys())
        st.header("My Uploaded Files")
        # Filter files for the current user (see user isolation in step 2)
        current_user = st.session_state.get("username", "unknown")
        available_files = list({
            md["filename"] 
            for md in metadata_store 
            if md.get("owner") == current_user or current_user in md.get("shared_with", [])
        })
        if available_files:
            if available_files:
                for i, fname in enumerate(available_files):
                    col1, col2 = st.columns([0.7, 0.3])
                    with col1:
                        st.write(fname)
                    with col2:
                        if st.button("Delete", key=f"del_{fname}_{i}"):
                            delete_file(fname)
        else:
            st.sidebar.info("No files uploaded yet.")

        # In the File Manager section:
        st.sidebar.header("Share a File")
        file_to_share = st.sidebar.selectbox("Select a file to share", available_files)
        share_with = st.sidebar.multiselect("Select user(s) to share with", options=available_usernames)

        if st.sidebar.button("Share File"):
            # Update metadata for the file if the current user is the owner
            for md in metadata_store:
                if md["filename"] == file_to_share and md.get("owner") == current_user:
                    md.setdefault("shared_with", []).extend(share_with)
                    md["shared_with"] = list(set(md["shared_with"]))  # Remove duplicates
                    st.success(f"Shared {file_to_share} with {', '.join(share_with)}")
            save_index_and_metadata()

    elif option == "Query Documents":
        st.header("Query Documents")
        st.sidebar.header("Settings")
        llm_model = st.sidebar.selectbox("Choose Your Model", ["Claude 3"])

        # "New Chat" button resets conversation and state.
        if st.sidebar.button("New Chat"):
            st.session_state.current_conversation_id = None  # Clear any old conversation ID
            st.session_state.current_conversation = None
            st.session_state.messages = []
            st.session_state.sources = []
            # Clear previous file/page selections too.
            st.session_state.selected_files = []
            st.session_state.selected_page_ranges = {}
            st.success("Started a new conversation.")
        web_search = st.sidebar.toggle("Enable Web Search")
        top_k = st.sidebar.slider("Select Top-K Results", min_value=1, max_value=100, value=50, step=1)

        # File and Page Range Selection
        # available_files = list(set([metadata['filename'] for metadata in metadata_store]))
        current_user = st.session_state.get("username", "unknown")
        # Only include files where the owner is the current user or shared with the user.
        available_files = list({
            md["filename"] 
            for md in metadata_store 
            if md.get("owner") == current_user or current_user in md.get("shared_with", [])
        })


        if available_files:
            # Use multiselect and store the selection in session state.
            st.session_state.selected_files = st.multiselect(
                "Select files (leave blank to search all):",
                available_files,
                default=st.session_state.selected_files
            )
            if len(st.session_state.selected_files) > 4:
                st.warning("For best results, select a maximum of 4 files.")
                # return

            page_ranges = get_page_range_for_files(st.session_state.selected_files)
            selected_page_ranges = {}
            # For each file, show page range inputs and store values in session state.
            for file in st.session_state.selected_files:
                min_page, max_page = page_ranges[file]
                col1, col2 = st.sidebar.columns(2)
                with col1:
                    start_page = st.number_input(
                        f"Start page for {file}",
                        min_value=min_page,
                        max_value=max_page,
                        value=page_ranges[file][0],
                        key=f"start_{file}"
                    )
                with col2:
                    end_page = st.number_input(
                        f"End page for {file}",
                        min_value=min_page,
                        max_value=max_page,
                        value=page_ranges[file][1],
                        key=f"end_{file}"
                    )
                selected_page_ranges[file] = (start_page, end_page)
            st.session_state.selected_page_ranges = selected_page_ranges

        st.sidebar.header("Previous Conversations")
        user_conversations = st.session_state.chat_history.get(current_user, [])

        # NEW CODE: Sort the entire list by timestamp
        unique_conversations = sorted(
            user_conversations, 
            key=lambda x: x.get("timestamp", ""), 
            reverse=True
        )

        for conv in unique_conversations:
            # Use the conversation's timestamp as a unique identifier.
            conv_id = conv.get("timestamp")
            default_label = conv.get("label") or conv.get('messages', [{}])[0].get("content", "")[:50]
            
            # Check if this conversation is in rename mode.
            if st.session_state.get("rename_mode") == conv_id:
                new_label = st.sidebar.text_input("Rename Conversation", value=default_label, key=f"rename_input_{conv_id}")
                if st.sidebar.button("Save", key=f"save_rename_{conv_id}"):
                    user = st.session_state.username
                    if user in st.session_state.chat_history:
                        for idx, stored_conv in enumerate(st.session_state.chat_history[user]):
                            if stored_conv.get("timestamp") == conv_id:
                                # Only update the label
                                st.session_state.chat_history[user][idx]["label"] = new_label
                                break 
                    save_chat_history(st.session_state.chat_history)
                    st.session_state["rename_mode"] = None
                    st.sidebar.success("✅ Conversation renamed successfully!")
                    st.rerun()
            else:
                col1, col2, col3, col4 = st.sidebar.columns([0.5, 0.2, 0.2, 0.1])
                if col1.button(default_label, key=f"load_{conv_id}"):
                    st.session_state.current_conversation = conv
                    st.session_state.messages = conv.get('messages', [])
                    st.session_state.selected_files = conv.get('files', [])
                    st.session_state.selected_page_ranges = conv.get('page_ranges', {})
                    st.session_state.current_conversation_id = conv_id
                    st.rerun()
                if col2.button("✏️", key=f"rename_button_{conv_id}"):
                    st.session_state["rename_mode"] = conv_id
                    st.rerun()
                if col3.button("🗑️", key=f"delete_{conv_id}"):
                    st.session_state["confirm_delete_conv"] = conv
                    st.rerun()
                if col4.button("📤", key=f"share_chat_{conv_id}"):
                    st.session_state["share_chat_conv"] = conv
                    st.session_state["share_chat_conv_id"] = conv_id
                    st.session_state["share_chat_mode"] = True
                    st.rerun()

        # Insert the share-chat snippet below the conversation list:
        if st.session_state.get("share_chat_mode"):
            st.header("Share Chat Conversation")
            share_chat_with = st.multiselect("Select user(s) to share with", options=available_usernames)
            if st.button("Confirm Share Chat"):
                chat_to_share = st.session_state["share_chat_conv"]
                # For each target user, append a deep copy of the conversation
                for user in share_chat_with:
                    if user in st.session_state.chat_history:
                        st.session_state.chat_history[user].append(copy.deepcopy(chat_to_share))
                    else:
                        st.session_state.chat_history[user] = [copy.deepcopy(chat_to_share)]
                save_chat_history(st.session_state.chat_history)
                st.success("Chat conversation shared successfully!")
                # Reset share mode variables
                st.session_state["share_chat_mode"] = False
                st.session_state.pop("share_chat_conv", None)
                st.session_state.pop("share_chat_conv_id", None)
                st.rerun()

        # If a conversation is marked for deletion, confirm deletion.
        if "confirm_delete_conv" in st.session_state:
            chat_name = (
                st.session_state["confirm_delete_conv"].get("label")
                or st.session_state["confirm_delete_conv"].get('messages', [{}])[0].get("content", "")[:50]
            )
            st.warning(f"Are you sure you want to delete '{chat_name}' conversation?")
            ccol1, ccol2 = st.columns(2)
            with ccol1:
                if st.button("Confirm Delete"):
                    user = st.session_state.username
                    if user in st.session_state.chat_history:
                        try:
                            st.session_state.chat_history[user].remove(st.session_state["confirm_delete_conv"])
                        except ValueError:
                            pass  # Conversation already removed.
                        save_chat_history(st.session_state.chat_history)
                    del st.session_state["confirm_delete_conv"]
                    st.sidebar.success("Conversation deleted!")
                    st.rerun()
            with ccol2:
                if st.button("Cancel"):
                    del st.session_state["confirm_delete_conv"]
                    st.sidebar.info("Deletion canceled.")
                    st.rerun()

        # --- Display chat messages with share options ---
        for idx, message in enumerate(st.session_state.messages):
            with st.chat_message(message["role"]):
                # Show role and time if available
                msg_time = message.get("time", "")
                role_title = message["role"].title()  # "User" / "Assistant"
                st.markdown(f"**`[{role_title} @ {msg_time}]`**\n\n{message['content']}")

                with st.expander("Show Copyable Text"):
                    st.code(message["content"], language="text")

        # --- New User Input using text_area ---
        user_message = user_input()
        if user_message:
            ist_timezone = pytz.timezone("Asia/Kolkata")
            timestamp_now = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")

            # 1) Append the user's message to st.session_state.messages
            st.session_state.messages.append({
                "role": "user",
                "content": user_message,
                "time": timestamp_now
            })

            # Display the user message in chat
            with st.chat_message("user"):
                st.markdown(user_message)

            # Prepare the last few messages for context.
            last_messages = st.session_state.messages[-5:] if len(st.session_state.messages) >= 5 else st.session_state.messages


            with st.spinner("Searching documents..."):
                st.markdown("**While you wait, Feel free to Refer to the Original Documents**")

                for file_key in st.session_state.selected_files:
                    # Generate the pre-signed URL for each file
                    preview_url = get_presigned_url(file_key)
                    # Create a clickable markdown link; clicking it will open the file in a new tab
                    st.markdown(f"[**{file_key}**]({preview_url})", unsafe_allow_html=True)


                answer = query_documents_with_page_range(
                    st.session_state.selected_files, 
                    st.session_state.selected_page_ranges, 
                    user_message,
                    top_k,
                    last_messages,
                    web_search
                )

                st.session_state.sources.append({
                    "answer": answer
                })


            ist_timezone = pytz.timezone("Asia/Kolkata")

            # After you get the LLM's answer, you append the "assistant" message:
            assistant_answer = answer  # from your retrieval
            current_time = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages.append({
                "role": "assistant",
                "content": assistant_answer,
                "time": current_time
            })

            # Show the assistant response in the UI
            with st.chat_message("assistant"):
                st.markdown(assistant_answer)

            user = st.session_state.username
            if user not in st.session_state.chat_history:
                st.session_state.chat_history[user] = []

            # Check if this is a continuation of an old conversation
            original_conv_id = st.session_state.get("original_conversation_id")

            if original_conv_id:
                # This is a continuation of an old conversation
                # 1. Find the original conversation to get its label
                original_conv = None
                original_label = "Conversation"
                for conv in st.session_state.chat_history[user]:
                    if conv.get("timestamp") == original_conv_id:
                        original_conv = conv
                        original_label = conv.get("label", original_label)
                        break
            
                # 2. Create a new conversation with the old label
                new_conversation_id = str(uuid.uuid4())
                new_conversation = {
                    "conversation_id": new_conversation_id,
                    "label": original_label,
                    "timestamp": current_time,
                    "messages": st.session_state.messages,
                    "files": st.session_state.selected_files,
                    "page_ranges": st.session_state.selected_page_ranges
                }

                # 3. Remove the old conversation
                if original_conv:
                    try:
                        st.session_state.chat_history[user].remove(original_conv)
                    except ValueError:
                        pass

            # 4. Add the new conversation
                st.session_state.chat_history[user].append(new_conversation)

            # 5. Update session state
                st.session_state.current_conversation_id = new_conversation_id
                st.session_state.original_conversation_id = None  # Clear the original ID
                
            else:
                # Normal new conversation flow
                current_conversation_id = st.session_state.get("current_conversation_id")
                
                if current_conversation_id:
                    # Update existing conversation
                    for conv in st.session_state.chat_history[user]:
                        if conv.get("conversation_id") == current_conversation_id:
                            conv["messages"] = st.session_state.messages
                            conv["timestamp"] = current_time
                            conv["files"] = st.session_state.selected_files
                            conv["page_ranges"] = st.session_state.selected_page_ranges
                            break
                else:
                    # Create new conversation
                    new_conversation_id = str(uuid.uuid4())
                    st.session_state.current_conversation_id = new_conversation_id
                    
                    new_conversation = {
                        "conversation_id": new_conversation_id,
                        "label": user_message[:50],
                        "timestamp": current_time,
                        "messages": st.session_state.messages,
                        "files": st.session_state.selected_files,
                        "page_ranges": st.session_state.selected_page_ranges
                    }
                    
                    st.session_state.chat_history[user].append(new_conversation)

            # Sort and save
            st.session_state.chat_history[user] = sorted(
                st.session_state.chat_history[user],
                key=lambda x: x.get("timestamp", ""),
                reverse=True
            )

            save_chat_history(st.session_state.chat_history)
            st.rerun()

    elif option == "Usage Monitoring":
        st.header("Usage Monitoring")

        # 1. Gather all usernames from your USERS dictionary
        all_usernames = list(USERS.keys())

        # 2. Add a dropdown in the sidebar for period selection
        period_options = ["Last 3 Days", "Last 7 Days", "Last 14 Days", "Last 1 Month", "Last 3 Months"]
        selected_period = st.sidebar.selectbox("Select Period", period_options, index=3)  # default: Last 1 Month
        
        # Map each period to the corresponding number of days
        period_days = {
            "Last 3 Days": 3,
            "Last 7 Days": 7,
            "Last 14 Days": 14,
            "Last 1 Month": 30,
            "Last 3 Months": 90
        }
        selected_days = period_days[selected_period]
        
        # Display the subheader with the chosen period
        st.subheader(f"Usage Monitoring - {selected_period}")

        # 3. Add a multi-select for user filtering
        selected_users = st.sidebar.multiselect(
            "Select User(s) to Display",
            options=all_usernames,
            default=all_usernames  # By default, all are selected
        )

        # 4. Load the chat history (keys are usernames)
        chat_history = load_chat_history()

        # Build a list of usage records from chat_history
        # Each record is { 'user': ..., 'timestamp': datetime_object }
        records = []
        for user, conversations in chat_history.items():
            for conv in conversations:
                timestamp_str = conv.get("timestamp")
                if timestamp_str:
                    try:
                        ts = datetime.strptime(timestamp_str, "%Y-%m-%d %H:%M:%S")
                        records.append({"user": user, "timestamp": ts})
                    except Exception as e:
                        st.warning(f"Timestamp format error for user {user}: {e}")

        # 5. Check if we have any usage records at all
        if not records:
            st.info("No usage data available.")
            return

        # Convert records to a DataFrame
        df = pd.DataFrame(records)

        # 6. Filter by the selected date range
        today = datetime.today()
        start_date = today - timedelta(days=selected_days)
        df_period = df[df["timestamp"] >= start_date]

        # 7. Filter by the selected users from the multi-select
        df_period = df_period[df_period["user"].isin(selected_users)]

        # If no data after both filters, exit
        if df_period.empty:
            st.info("No usage data found for the chosen date range and user(s).")
            return

        # --- BAR CHART: Total Queries per User ---
        user_counts = df_period.groupby("user").size().reset_index(name="queries")
        bar_fig = px.bar(
            user_counts,
            x="user",
            y="queries",
            title=f"Total Queries per User ({selected_period})",
            labels={"user": "User", "queries": "Number of Queries"}
        )
        st.plotly_chart(bar_fig, use_container_width=True)

        # --- LINE CHART: Day-wise Queries per User with Moving Average ---
        # Create a 'date' column (just the date part)
        df_period["date"] = df_period["timestamp"].dt.date

        # Count queries per user per day
        daily_counts = df_period.groupby(["user", "date"]).size().reset_index(name="queries")

        # Create a complete date range for the selected period
        date_range = pd.date_range(start=start_date.date(), end=today.date())
        all_users_in_data = daily_counts["user"].unique()
        complete_data = []

        # Build a day-by-day dataset with possible missing dates filled in
        for user in all_users_in_data:
            user_df = daily_counts[daily_counts["user"] == user].copy()
            user_df.set_index("date", inplace=True)
            # Reindex to fill missing days with 0
            user_df = user_df.reindex(date_range, fill_value=0)
            user_df = user_df.rename_axis("date").reset_index()
            user_df["user"] = user
            
            # Calculate a moving average based on 'selected_days'
            user_df["moving_avg"] = user_df["queries"].rolling(window=selected_days, min_periods=1).mean()
            complete_data.append(user_df)

        # Combine each user's timeseries data
        daily_all = pd.concat(complete_data, ignore_index=True)

        # Create the line graph showing daily queries
        line_fig = px.line(
            daily_all,
            x="date",
            y="queries",
            color="user",
            title=f"Daily Queries per User ({selected_period})",
            labels={"date": "Date", "queries": "Number of Queries"}
        )

        # Add a separate trace for the rolling average of each user
        for user in all_users_in_data:
            user_data = daily_all[daily_all["user"] == user]
            line_fig.add_trace(
                go.Scatter(
                    x=user_data["date"],
                    y=user_data["moving_avg"],
                    mode="lines",
                    name=f"{user} - {selected_period} MA"
                )
            )

        st.plotly_chart(line_fig, use_container_width=True)

    else:
        st.warning("No files available in the index. Please upload Documents to populate the index.")

if __name__ == "__main__":
    main()
